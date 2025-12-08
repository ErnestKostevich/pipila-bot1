#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
🔽 PIPILA v8.5 FINAL - Document Processor with BATCH MODE
Creates ChromaDB from documents on each deploy
"""

import os
import sys
import urllib.request
import zipfile
import shutil
import time
from pathlib import Path

# ChromaDB path (local folder, recreated each deploy)
CHROMA_PATH = "./chroma_db"

def log(msg):
    print(f"[PROCESSOR] {msg}", flush=True)
    sys.stdout.flush()

def download_documents():
    """Download documents ZIP from GitHub Releases"""
    
    github_url = "https://github.com/ErnestKostevich/pipila-bot1/releases/download/v8.2/Fuentes.de.informacion.RAG-20251207T164947Z-3-001.zip"
    
    zip_path = "/tmp/documents.zip"
    extract_dir = "/tmp/documents"
    
    log("=" * 70)
    log("🔽 Downloading Documents from GitHub")
    log("=" * 70)
    
    # Clean old folders
    for path in [zip_path, extract_dir, CHROMA_PATH]:
        if os.path.exists(path):
            try:
                if os.path.isdir(path):
                    shutil.rmtree(path)
                else:
                    os.remove(path)
            except:
                pass
    
    log(f"📥 Downloading...")
    start_time = time.time()
    
    try:
        req = urllib.request.Request(github_url, headers={'User-Agent': 'Mozilla/5.0'})
        
        with urllib.request.urlopen(req, timeout=600) as response:
            with open(zip_path, 'wb') as out_file:
                total_size = int(response.headers.get('content-length', 0))
                log(f"📦 Size: {total_size / (1024*1024):.0f} MB")
                
                downloaded = 0
                chunk_size = 131072
                last_report = 0
                
                while True:
                    chunk = response.read(chunk_size)
                    if not chunk:
                        break
                    out_file.write(chunk)
                    downloaded += len(chunk)
                    
                    mb = downloaded / (1024 * 1024)
                    if mb - last_report >= 100:
                        log(f"   {mb:.0f} MB...")
                        last_report = mb
        
        elapsed = time.time() - start_time
        log(f"✅ Downloaded in {elapsed:.0f}s")
        
    except Exception as e:
        log(f"❌ Download FAILED: {e}")
        sys.exit(1)
    
    # Extract
    log("📦 Extracting...")
    try:
        os.makedirs(extract_dir, exist_ok=True)
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_dir)
        os.remove(zip_path)
        log("✅ Extracted")
    except Exception as e:
        log(f"❌ Extract FAILED: {e}")
        sys.exit(1)
    
    return extract_dir

def find_documents(base_dir):
    """Find all processable documents"""
    supported = {'.pdf', '.docx', '.pptx', '.txt'}
    documents = []
    
    for root, dirs, files in os.walk(base_dir):
        dirs[:] = [d for d in dirs if not d.startswith('.') and not d.startswith('~')]
        for file in files:
            if file.startswith('~') or file.startswith('.'):
                continue
            if Path(file).suffix.lower() in supported:
                documents.append(os.path.join(root, file))
    
    return documents

def extract_text_from_pdf(file_path):
    try:
        import PyPDF2
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                try:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
                except:
                    continue
            return text.strip()
    except:
        return ""

def extract_text_from_docx(file_path):
    try:
        import docx
        doc = docx.Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs if p.text]).strip()
    except:
        return ""

def extract_text_from_pptx(file_path):
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts).strip()
    except:
        return ""

def extract_text_from_txt(file_path):
    for enc in ['utf-8', 'latin-1', 'cp1252']:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                return f.read().strip()
        except:
            continue
    return ""

def extract_text(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_text_from_docx(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext == '.txt':
        return extract_text_from_txt(file_path)
    return ""

def chunk_text(text, chunk_size=1000, overlap=200):
    if not text or len(text) < 100:
        return []
    
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if len(chunk) > 50:
            chunks.append(chunk)
        start = end - overlap
        if start >= len(text):
            break
    return chunks

def create_chromadb(documents_dir):
    """Create ChromaDB with batch processing - THE WORKING VERSION"""
    import chromadb
    
    log("")
    log("=" * 70)
    log("🔧 Creating ChromaDB (BATCH MODE)")
    log("=" * 70)
    
    # Ensure directory exists
    os.makedirs(CHROMA_PATH, exist_ok=True)
    
    # Find documents
    log("📂 Scanning...")
    doc_files = find_documents(documents_dir)
    log(f"   Found {len(doc_files)} documents")
    
    if not doc_files:
        log("❌ No documents found!")
        sys.exit(1)
    
    # PHASE 1: Extract all text
    log("")
    log("📄 PHASE 1: Extracting text...")
    start_time = time.time()
    
    all_chunks = []
    all_metadatas = []
    all_ids = []
    
    processed = 0
    skipped = 0
    
    for i, doc_path in enumerate(doc_files):
        filename = os.path.basename(doc_path)
        
        if (i + 1) % 10 == 0:
            elapsed = time.time() - start_time
            rate = (i + 1) / elapsed if elapsed > 0 else 0
            eta = (len(doc_files) - i - 1) / rate / 60 if rate > 0 else 0
            log(f"   [{i+1}/{len(doc_files)}] {len(all_chunks)} chunks | ETA: {eta:.1f}min")
        
        try:
            text = extract_text(doc_path)
            
            if not text or len(text) < 100:
                skipped += 1
                continue
            
            chunks = chunk_text(text)
            if not chunks:
                skipped += 1
                continue
            
            for j, chunk in enumerate(chunks):
                all_chunks.append(chunk)
                all_metadatas.append({
                    "source": filename,
                    "chunk": j,
                    "total_chunks": len(chunks)
                })
                all_ids.append(f"doc_{processed}_{j}")
            
            processed += 1
            
        except Exception as e:
            skipped += 1
    
    phase1_time = time.time() - start_time
    log(f"✅ Phase 1: {len(all_chunks)} chunks in {phase1_time:.1f}s")
    
    # PHASE 2: Batch add to ChromaDB
    log("")
    log("🗄️ PHASE 2: Creating embeddings...")
    
    phase2_start = time.time()
    
    client = chromadb.PersistentClient(path=CHROMA_PATH)
    try:
        client.delete_collection("pipila_documents")
    except:
        pass
    
    collection = client.create_collection(
        name="pipila_documents",
        metadata={"hnsw:space": "cosine"}
    )
    
    # Add in batches of 500
    BATCH_SIZE = 500
    total_batches = (len(all_chunks) + BATCH_SIZE - 1) // BATCH_SIZE
    
    for batch_num in range(total_batches):
        start_idx = batch_num * BATCH_SIZE
        end_idx = min(start_idx + BATCH_SIZE, len(all_chunks))
        
        batch_chunks = all_chunks[start_idx:end_idx]
        batch_metas = all_metadatas[start_idx:end_idx]
        batch_ids = all_ids[start_idx:end_idx]
        
        elapsed = time.time() - phase2_start
        if batch_num > 0:
            rate = batch_num / elapsed
            eta = (total_batches - batch_num) / rate / 60
            log(f"   Batch {batch_num + 1}/{total_batches} | ETA: {eta:.0f}min")
        else:
            log(f"   Batch {batch_num + 1}/{total_batches}...")
        
        collection.add(
            documents=batch_chunks,
            metadatas=batch_metas,
            ids=batch_ids
        )
    
    phase2_time = time.time() - phase2_start
    total_time = time.time() - start_time
    
    log("")
    log("=" * 70)
    log("✅ ChromaDB Created!")
    log("=" * 70)
    log(f"📄 Documents: {processed}")
    log(f"📊 Chunks: {len(all_chunks)}")
    log(f"📁 Saved to: {CHROMA_PATH}")
    log(f"⏱️ Time: {total_time/60:.1f} minutes")
    log("=" * 70)
    
    return len(all_chunks)

def main():
    log("=" * 70)
    log("🚀 PIPILA v8.5 FINAL - BATCH MODE")
    log("=" * 70)
    log("")
    
    total_start = time.time()
    
    # Download documents
    documents_dir = download_documents()
    
    # Install pptx if needed
    try:
        from pptx import Presentation
    except ImportError:
        log("📦 Installing python-pptx...")
        os.system("pip install python-pptx --quiet --break-system-packages 2>/dev/null || pip install python-pptx --quiet")
    
    # Create ChromaDB with BATCH MODE
    total_chunks = create_chromadb(documents_dir)
    
    # Cleanup
    log("🧹 Cleaning up...")
    try:
        shutil.rmtree(documents_dir)
    except:
        pass
    
    total_time = time.time() - total_start
    
    log("")
    log("=" * 70)
    log("🎉 ALL DONE!")
    log("=" * 70)
    log(f"📊 ChromaDB: {total_chunks} chunks")
    log(f"📁 Path: {CHROMA_PATH}")
    log(f"⏱️ Total: {total_time/60:.1f} minutes")
    log("=" * 70)
    log("")
    log("✅ Starting bot...")
    
    return 0

if __name__ == "__main__":
    try:
        sys.exit(main())
    except KeyboardInterrupt:
        log("⚠️ Interrupted")
        sys.exit(1)
    except Exception as e:
        log(f"💥 ERROR: {e}")
        import traceback
        log(traceback.format_exc())
        sys.exit(1)
